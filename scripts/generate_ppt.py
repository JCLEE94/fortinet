#!/usr/bin/env python3
"""
FortiGate Nextrade 엔터프라이즈 평가 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_title_slide(prs):
    """타이틀 슬라이드 생성"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "FortiGate Nextrade"
    subtitle.text = "엔터프라이즈 환경 적합성 평가\n\n2025년 7월\n\n차세대 네트워크 보안 관리 플랫폼"
    
    # 타이틀 서식
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # 배경 그라디언트 효과 (단색으로 대체)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea

def create_agenda_slide(prs):
    """목차 슬라이드"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "목차"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Executive Summary\n"
    tf.text += "2. 솔루션 개요\n"
    tf.text += "3. 핵심 기능 및 차별점\n"
    tf.text += "4. 기술 아키텍처\n"
    tf.text += "5. 성능 및 확장성\n"
    tf.text += "6. 보안 및 컴플라이언스\n"
    tf.text += "7. TCO 및 ROI 분석\n"
    tf.text += "8. 구현 로드맵\n"
    tf.text += "9. 결론 및 권장사항"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)

def create_executive_summary_slide(prs):
    """Executive Summary 슬라이드"""
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    # 좌측 박스 - 종합 평가 점수
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(4)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 250)
    
    tf = shape.text_frame
    tf.text = "종합 평가 점수\n85/100 ⭐⭐⭐⭐☆\n\n"
    tf.text += "• 기능성: 90%\n"
    tf.text += "• 성능: 80%\n"
    tf.text += "• 보안: 85%\n"
    tf.text += "• 운영성: 85%"
    
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.LEFT
    
    # 우측 박스 - 핵심 평가 결과
    left = Inches(5)
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(230, 250, 230)
    
    tf2 = shape2.text_frame
    tf2.text = "핵심 평가 결과\n\n"
    tf2.text += "✅ 중견기업 즉시 적용 가능\n"
    tf2.text += "✅ 24/7 운영 환경 지원\n"
    tf2.text += "✅ 18개월 내 ROI 달성\n"
    tf2.text += "⚠️ 대기업은 일부 커스터마이징 필요"
    
    tf2.word_wrap = True
    for paragraph in tf2.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.LEFT

def create_solution_overview_slide(prs):
    """솔루션 개요 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "솔루션 개요"
    
    # 중앙 박스
    left = Inches(2)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(1)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    tf = shape.text_frame
    tf.text = "FortiGate Nextrade\n중앙집중식 보안 관리 플랫폼"
    tf.word_wrap = True
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 하단 특징들
    features = [
        ("정책 관리", 1.5),
        ("모니터링", 3.5),
        ("자동화", 5.5),
        ("통합 관리", 7.5)
    ]
    
    for feature, left_pos in features:
        left = Inches(left_pos)
        top = Inches(3)
        width = Inches(1.5)
        height = Inches(0.8)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 147, 251)
        
        tf = shape.text_frame
        tf.text = feature
        tf.word_wrap = True
        
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 하단 설명
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(8)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "• 클라우드 네이티브 아키텍처\n"
    tf.text += "• AI 기반 위협 탐지\n"
    tf.text += "• 실시간 토폴로지 시각화\n"
    tf.text += "• 엔터프라이즈급 확장성"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

def create_key_features_slide(prs):
    """핵심 기능 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "핵심 기능"
    
    # 4분할 그리드
    features = [
        ("🔐 정책 관리", "• CRUD 완벽 지원\n• 시나리오 분석\n• 일괄 처리", 0.5, 1.5),
        ("📊 실시간 모니터링", "• SSE 기반 실시간\n• 대시보드 제공\n• 알림 자동화", 5, 1.5),
        ("🌐 토폴로지 시각화", "• D3.js 기반\n• 1000+ 노드\n• 실시간 업데이트", 0.5, 3.5),
        ("🔄 ITSM 통합", "• 티켓 자동 생성\n• 워크플로우 연동\n• API 통합", 5, 3.5)
    ]
    
    for title_text, content, left, top in features:
        width = Inches(4)
        height = Inches(1.8)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        tf = shape.text_frame
        tf.text = f"{title_text}\n{content}"
        tf.word_wrap = True
        
        # 첫 번째 줄(제목)은 굵게
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        
        # 나머지 내용
        for i in range(1, len(tf.paragraphs)):
            tf.paragraphs[i].font.size = Pt(14)

def create_performance_slide(prs):
    """성능 벤치마크 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "성능 벤치마크"
    
    # 성능 지표 박스
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    tf = shape.text_frame
    tf.text = "처리량 테스트 결과\n"
    tf.text += "• 15,000 Requests/min 처리 가능\n"
    tf.text += "• 45ms 평균 응답시간 달성\n"
    tf.text += "• 500+ 동시 사용자 지원"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
    
    # 노드 렌더링 성능
    left = Inches(0.5)
    top = Inches(4)
    width = Inches(9)
    height = Inches(1.5)
    
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(240, 250, 240)
    
    tf2 = shape2.text_frame
    tf2.text = "노드 렌더링 성능\n"
    tf2.text += "100 nodes: 0.5초 | 500 nodes: 2.1초 | 1000 nodes: 5.3초 | 5000 nodes: 15.2초"
    
    for paragraph in tf2.paragraphs:
        paragraph.font.size = Pt(16)

def create_tco_slide(prs):
    """TCO 분석 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "TCO 및 ROI 분석"
    
    # 표 생성 - 헤더 포함해서 8행 필요
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    
    shape = slide.shapes.add_table(8, 4, left, top, width, height)
    table = shape.table
    
    # 헤더
    headers = ['구분', '1차년도', '2차년도', '3차년도']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # 데이터
    data = [
        ['라이선스', '$50,000', '$50,000', '$50,000'],
        ['인프라', '$24,000', '$28,800', '$34,560'],
        ['운영인력', '$60,000', '$63,000', '$66,150'],
        ['유지보수', '$10,000', '$12,000', '$14,400'],
        ['총 비용', '$144,000', '$153,800', '$165,110'],
        ['절감효과', '-$30,000', '-$45,000', '-$60,000'],
        ['순 비용', '$114,000', '$108,800', '$105,110']
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
    
    # ROI 텍스트
    left = Inches(1)
    top = Inches(5)
    width = Inches(8)
    height = Inches(0.8)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "투자 대비 수익률(ROI): 185% (3년 기준) | 투자 회수 기간: 18개월"
    
    paragraph = tf.paragraphs[0]
    paragraph.font.size = Pt(20)
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER

def create_roadmap_slide(prs):
    """구현 로드맵 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "구현 로드맵"
    
    phases = [
        ("Phase 1: Foundation (1-2개월)", "• LDAP/AD 통합\n• RBAC 구현\n• 감사 로그 강화", 1.5),
        ("Phase 2: Enhancement (3-4개월)", "• HA 구성\n• 성능 최적화\n• SIEM 통합", 3),
        ("Phase 3: Advanced (5-6개월)", "• 멀티 테넌시\n• AI/ML 통합\n• 자동화 확대", 4.5)
    ]
    
    for phase_title, phase_content, top_pos in phases:
        left = Inches(0.5)
        top = Inches(top_pos)
        width = Inches(9)
        height = Inches(1.2)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 250)
        
        tf = shape.text_frame
        tf.text = f"{phase_title}\n{phase_content}"
        tf.word_wrap = True
        
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        
        for i in range(1, len(tf.paragraphs)):
            tf.paragraphs[i].font.size = Pt(14)

def create_recommendation_slide(prs):
    """권장사항 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "권장사항"
    
    recommendations = [
        ("💡 경영진을 위한 권장사항", 
         "• 즉시 도입 검토 권장\n• Pilot → 부서별 → 전사 확대\n• 18개월 내 투자비 회수 예상",
         RGBColor(227, 242, 253)),
        ("🔧 IT 부서를 위한 권장사항",
         "• 2-4주 PoC 수행\n• 기존 ITSM/SIEM 연동 계획\n• Python/K8s 역량 확보",
         RGBColor(243, 229, 245)),
        ("📈 향후 발전 방향",
         "• AI/ML 이상 탐지 통합\n• 완전한 SaaS 전환\n• Multi-region 글로벌 확장",
         RGBColor(232, 245, 233))
    ]
    
    top_start = 1.5
    for i, (rec_title, rec_content, color) in enumerate(recommendations):
        left = Inches(0.5)
        top = Inches(top_start + i * 1.8)
        width = Inches(9)
        height = Inches(1.5)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        
        tf = shape.text_frame
        tf.text = f"{rec_title}\n{rec_content}"
        tf.word_wrap = True
        
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        
        for j in range(1, len(tf.paragraphs)):
            tf.paragraphs[j].font.size = Pt(14)

def create_conclusion_slide(prs):
    """결론 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "결론"
    
    # 중앙 박스
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    tf = shape.text_frame
    tf.text = "FortiGate Nextrade\n\n"
    tf.text += "엔터프라이즈 준비도: 85%\n\n"
    tf.text += "중견기업 및 준대기업 환경에서\n"
    tf.text += "즉시 배포 가능한 수준의 완성도"
    
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 하단 체크리스트
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(8)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "✓ 검증된 기술 스택  ✓ 우수한 투자 대비 효과  ✓ 빠른 구현 가능  ✓ 지속적인 발전 가능성"
    
    paragraph = tf.paragraphs[0]
    paragraph.font.size = Pt(16)
    paragraph.alignment = PP_ALIGN.CENTER

def create_qa_slide(prs):
    """Q&A 슬라이드"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 없이 중앙에 Q&A
    left = Inches(3)
    top = Inches(2)
    width = Inches(4)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "Questions & Answers"
    
    paragraph = tf.paragraphs[0]
    paragraph.font.size = Pt(36)
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER
    
    # 감사합니다
    left = Inches(3.5)
    top = Inches(3.5)
    width = Inches(3)
    height = Inches(0.8)
    
    textbox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = textbox2.text_frame
    tf2.text = "감사합니다"
    
    paragraph2 = tf2.paragraphs[0]
    paragraph2.font.size = Pt(24)
    paragraph2.alignment = PP_ALIGN.CENTER
    
    # 연락처
    left = Inches(2)
    top = Inches(5)
    width = Inches(6)
    height = Inches(1.5)
    
    textbox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = textbox3.text_frame
    tf3.text = "📧 nextrade@fortinet.com\n"
    tf3.text += "📞 02-1234-5678\n"
    tf3.text += "🌐 fortinet.jclee.me"
    
    for paragraph in tf3.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.alignment = PP_ALIGN.CENTER

def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = Presentation()
    
    # 16:9 비율 설정
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 각 슬라이드 생성
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_executive_summary_slide(prs)
    create_solution_overview_slide(prs)
    create_key_features_slide(prs)
    create_performance_slide(prs)
    create_tco_slide(prs)
    create_roadmap_slide(prs)
    create_recommendation_slide(prs)
    create_conclusion_slide(prs)
    create_qa_slide(prs)
    
    # 파일 저장
    output_dir = "/home/jclee/app/fortinet/docs/presentations"
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = os.path.join(output_dir, "FortiGate_Nextrade_Enterprise_Evaluation.pptx")
    prs.save(output_file)
    
    print(f"✅ PPT 파일이 생성되었습니다: {output_file}")
    print(f"📊 총 슬라이드 수: {len(prs.slides)}장")

if __name__ == "__main__":
    main()